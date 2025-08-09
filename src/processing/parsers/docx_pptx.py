# Parsers for Office documents (functional style)

from pathlib import Path
from typing import Dict, List
import zipfile

from ..normalize import normalize_newlines, utf8_decode_remove_bom, remove_control_chars
from ..io import ensure_output_dir, write_text_file
from .. import mapping

_HAS_DOCX = False
_HAS_PPTX = False
try:
    import docx as _docx_mod  # type: ignore

    _HAS_DOCX = True
except Exception:
    _HAS_DOCX = False

try:
    from pptx import Presentation as _Presentation  # type: ignore

    _HAS_PPTX = True
except Exception:
    _HAS_PPTX = False

def _ensure_output_dir(base_dir: Path) -> Path:
    # Backward-compatible wrapper calling shared I/O helper.
    return ensure_output_dir(base_dir)


def _write_text_file(path: Path, text: str) -> None:
    # Backward-compatible wrapper calling shared I/O helper.
    write_text_file(path, text)


def _extract_all_between(s: str, start: str, end: str) -> List[str]:
    # Simple scanning without regex
    results = []
    idx = 0
    while True:
        i = s.find(start, idx)
        if i == -1:
            break
        j = s.find(end, i + len(start))
        if j == -1:
            break
        results.append(s[i + len(start) : j])
        idx = j + len(end)
    return results


def parse_docx(src_path: Path, base_dir: Path) -> Dict[str, object]:
    p = Path(src_path)
    base = Path(base_dir)

    paragraphs: List[str] = []
    if _HAS_DOCX:
        try:
            doc = _docx_mod.Document(str(p))
            i = 0
            while i < len(doc.paragraphs):
                paragraphs.append(doc.paragraphs[i].text)
                i = i + 1
        except Exception:
            # Fallback to zip-based parsing
            with zipfile.ZipFile(p, "r") as z:
                xml_bytes = z.read("word/document.xml")
                xml = utf8_decode_remove_bom(xml_bytes)
                texts = _extract_all_between(xml, "<w:t>", "</w:t>")
                j = 0
                while j < len(texts):
                    paragraphs.append(texts[j])
                    j = j + 1
    else:
        with zipfile.ZipFile(p, "r") as z:
            xml_bytes = z.read("word/document.xml")
            xml = utf8_decode_remove_bom(xml_bytes)
            texts = _extract_all_between(xml, "<w:t>", "</w:t>")
            j = 0
            while j < len(texts):
                paragraphs.append(texts[j])
                j = j + 1

    # Join as lines and normalize
    text = "\n".join(paragraphs)
    text = normalize_newlines(text)
    text = remove_control_chars(text)

    out_dir = _ensure_output_dir(base)
    out_name = p.stem + ".txt"
    out_path = out_dir / out_name
    _write_text_file(out_path, text)

    map_entry = mapping.capture_paths(p, out_path)
    return {"out_path": str(out_path), "text": text, "mapping": map_entry}


def _slide_index_from_name(name: str) -> int:
    # name like 'ppt/slides/slideN.xml' → return N
    # Parse digits after last 'slide'
    idx = name.rfind("slide")
    if idx == -1:
        return 0
    i = idx + len("slide")
    n = 0
    found = False
    while i < len(name):
        ch = name[i]
        if ch >= "0" and ch <= "9":
            found = True
            n = n * 10 + (ord(ch) - ord("0"))
            i = i + 1
            continue
        break
    if not found:
        return 0
    return n


def parse_pptx(src_path: Path, base_dir: Path) -> Dict[str, object]:
    p = Path(src_path)
    base = Path(base_dir)

    slides_texts: List[str] = []
    if _HAS_PPTX:
        try:
            prs = _Presentation(str(p))
            si = 0
            while si < len(prs.slides):
                slide = prs.slides[si]
                chunks = []
                # Walk shapes and collect text runs
                sj = 0
                shapes = slide.shapes
                while sj < len(shapes):
                    shp = shapes[sj]
                    try:
                        if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                            tf = shp.text_frame
                            pi = 0
                            while pi < len(tf.paragraphs):
                                para = tf.paragraphs[pi]
                                ri = 0
                                while ri < len(para.runs):
                                    chunks.append(para.runs[ri].text)
                                    ri = ri + 1
                                pi = pi + 1
                    except Exception:
                        # Ignore shapes we cannot parse
                        pass
                    sj = sj + 1
                slide_text = "\n".join(chunks)
                slide_text = normalize_newlines(slide_text)
                slides_texts.append(slide_text.rstrip("\n"))
                si = si + 1
        except Exception:
            # Fallback to zip-based parsing
            with zipfile.ZipFile(p, "r") as z:
                names = []
                for name in z.namelist():
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                        names.append(name)
                ordered = []
                while len(names) > 0:
                    best_idx = 0
                    k = 1
                    while k < len(names):
                        if _slide_index_from_name(names[k]) < _slide_index_from_name(names[best_idx]):
                            best_idx = k
                        k = k + 1
                    ordered.append(names[best_idx])
                    del names[best_idx]
                si = 0
                while si < len(ordered):
                    name = ordered[si]
                    xml_bytes = z.read(name)
                    xml = utf8_decode_remove_bom(xml_bytes)
                    texts = _extract_all_between(xml, "<a:t>", "</a:t>")
                    slide_text = "\n".join(texts)
                    slide_text = normalize_newlines(slide_text)
                    slides_texts.append(slide_text.rstrip("\n"))
                    si = si + 1
    else:
        with zipfile.ZipFile(p, "r") as z:
            names = []
            for name in z.namelist():
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    names.append(name)
            ordered = []
            while len(names) > 0:
                best_idx = 0
                k = 1
                while k < len(names):
                    if _slide_index_from_name(names[k]) < _slide_index_from_name(names[best_idx]):
                        best_idx = k
                    k = k + 1
                ordered.append(names[best_idx])
                del names[best_idx]
            si = 0
            while si < len(ordered):
                name = ordered[si]
                xml_bytes = z.read(name)
                xml = utf8_decode_remove_bom(xml_bytes)
                texts = _extract_all_between(xml, "<a:t>", "</a:t>")
                slide_text = "\n".join(texts)
                slide_text = normalize_newlines(slide_text)
                slides_texts.append(slide_text.rstrip("\n"))
                si = si + 1

    # Join slides with separator line
    combined = "\n---\n".join(slides_texts)
    text = normalize_newlines(combined)
    text = remove_control_chars(text)

    out_dir = _ensure_output_dir(base)
    out_name = p.stem + ".txt"
    out_path = out_dir / out_name
    _write_text_file(out_path, text)

    map_entry = mapping.capture_paths(p, out_path)
    return {"out_path": str(out_path), "text": text, "mapping": map_entry}
